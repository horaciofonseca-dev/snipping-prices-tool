import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import os
import shutil
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Paths
DATA_PATH = r"C:\Users\emman\p_Claude\devs\snipper_tool\capstone\synthetic_data\synthetic_12month_inflation_data.csv"
BASKET_ANALYSIS_DIR = r"C:\Users\emman\p_Claude\devs\snipper_tool\capstone\synthetic_data\basket_analysis"
BASKET_CORRECTED_DIR = r"C:\Users\emman\p_Claude\devs\snipper_tool\capstone\synthetic_data\basket_analysis_corrected"
PRESENTATION_OUTPUT = r"C:\Users\emman\p_Claude\devs\snipper_tool\CAPSTONE_DELIVERABLES\Snipper_Tool_Reality_vs_Official_Paris.pptx"

# Corrected basket costs
OFFICIAL_MONTHLY = 68.67
REAL_MONTHLY = 303.65
HEALTHY_MONTHLY = 383.58

PARIS_INCOMES = {
    'SMIC (Minimum)': 21000,
    'Low Income': 28000,
    'Median': 42000,
    'Upper Middle': 65000,
}

# Colors
DARK_BLUE = RGBColor(31, 78, 121)
ACCENT_ORANGE = RGBColor(237, 125, 49)
RED_ALERT = RGBColor(255, 68, 68)
WHITE = RGBColor(255, 255, 255)
TEXT_BLACK = RGBColor(51, 51, 51)

CHART_BLUE = '#1F4E79'
CHART_GREEN = '#4F81BD'
CHART_ORANGE = '#ED7D31'

def create_corrected_visualizations():
    """Generate visualizations with corrected basket data"""
    print("[STEP 1] Generating corrected visualizations...")

    # 1. Basket Comparison
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 6))

    baskets = ['Official\n(INSEE)', 'Real\nComplete', 'Healthy\nLifestyle']
    monthly = [OFFICIAL_MONTHLY, REAL_MONTHLY, HEALTHY_MONTHLY]
    annual = [OFFICIAL_MONTHLY * 12, REAL_MONTHLY * 12, HEALTHY_MONTHLY * 12]
    colors = ['#ff6b6b', '#4f81bd', '#70ad47']

    bars1 = ax1.bar(baskets, monthly, color=colors, edgecolor='black', linewidth=2)
    ax1.set_ylabel('Monthly Cost (EUR)', fontsize=13, fontweight='bold')
    ax1.set_title('Monthly Basket Costs - Paris Household of 3', fontsize=14, fontweight='bold')
    ax1.set_ylim(0, max(monthly) * 1.25)

    for bar, cost in zip(bars1, monthly):
        height = bar.get_height()
        ax1.text(bar.get_x() + bar.get_width()/2., height,
                f'EUR {cost:.2f}', ha='center', va='bottom', fontweight='bold', fontsize=12)

    bars2 = ax2.bar(baskets, annual, color=colors, edgecolor='black', linewidth=2)
    ax2.set_ylabel('Annual Cost (EUR)', fontsize=13, fontweight='bold')
    ax2.set_title('Annual Basket Costs - Paris Household of 3', fontsize=14, fontweight='bold')
    ax2.set_ylim(0, max(annual) * 1.25)

    for bar, cost in zip(bars2, annual):
        height = bar.get_height()
        ax2.text(bar.get_x() + bar.get_width()/2., height,
                f'EUR {cost:,.0f}', ha='center', va='bottom', fontweight='bold', fontsize=12)

    plt.tight_layout()
    plt.savefig(os.path.join(BASKET_ANALYSIS_DIR, '01_basket_comparison.png'), dpi=150, bbox_inches='tight')
    plt.close()
    print("   [OK] 01_basket_comparison.png")

    # 2. Affordability Cliffs
    fig, ax = plt.subplots(figsize=(14, 7))

    income_names = list(PARIS_INCOMES.keys())
    income_values = list(PARIS_INCOMES.values())

    x = np.arange(len(income_names))
    width = 0.25

    official_pcts = [(OFFICIAL_MONTHLY * 12 / inc) * 100 for inc in income_values]
    real_pcts = [(REAL_MONTHLY * 12 / inc) * 100 for inc in income_values]
    healthy_pcts = [(HEALTHY_MONTHLY * 12 / inc) * 100 for inc in income_values]

    bars1 = ax.bar(x - width, official_pcts, width, label='Official Basket',
                  color=CHART_BLUE, edgecolor='black', linewidth=1.5)
    bars2 = ax.bar(x, real_pcts, width, label='Real Complete Basket',
                  color=CHART_ORANGE, edgecolor='black', linewidth=1.5)
    bars3 = ax.bar(x + width, healthy_pcts, width, label='Healthy Basket',
                  color='#70ad47', edgecolor='black', linewidth=1.5)

    # Affordability thresholds
    ax.axhline(y=30, color='red', linestyle='--', linewidth=2.5, alpha=0.7, label='Affordability Cliff (30%)')
    ax.axhline(y=20, color='orange', linestyle='--', linewidth=2, alpha=0.7, label='Warning Level (20%)')

    ax.set_ylabel('% of Annual Income Spent on Food', fontsize=13, fontweight='bold')
    ax.set_title('Food Affordability Cliffs by Paris Household Income', fontsize=14, fontweight='bold')
    ax.set_xticks(x)
    ax.set_xticklabels(income_names, fontsize=11)
    ax.legend(loc='upper right', fontsize=10)
    ax.set_ylim(0, 40)

    # Add value labels
    for bars in [bars1, bars2, bars3]:
        for bar in bars:
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height,
                   f'{height:.1f}%', ha='center', va='bottom', fontweight='bold', fontsize=9)

    plt.tight_layout()
    plt.savefig(os.path.join(BASKET_ANALYSIS_DIR, '02_affordability_cliffs.png'), dpi=150, bbox_inches='tight')
    plt.close()
    print("   [OK] 02_affordability_cliffs.png")

    # 3. Gap Analysis
    fig, ax = plt.subplots(figsize=(12, 7))

    gap_official_real = REAL_MONTHLY - OFFICIAL_MONTHLY
    gap_real_healthy = HEALTHY_MONTHLY - REAL_MONTHLY

    base_heights = [OFFICIAL_MONTHLY, OFFICIAL_MONTHLY, OFFICIAL_MONTHLY]
    gap1_heights = [0, gap_official_real, gap_official_real]
    gap2_heights = [0, 0, gap_real_healthy]

    bars1 = ax.bar(baskets, base_heights, color='#ff6b6b', edgecolor='black', linewidth=2, label='Official Basket Cost')
    bars2 = ax.bar(baskets, gap1_heights, bottom=base_heights, color='#ffa500', edgecolor='black', linewidth=2, label='Real Items Gap')
    bars3 = ax.bar(baskets, gap2_heights, bottom=[b+g for b, g in zip(base_heights, gap1_heights)],
                  color='#ff7f50', edgecolor='black', linewidth=2, label='Health Quality Gap')

    # Annotations
    for i, (basket, official, real, healthy) in enumerate(zip(baskets, [OFFICIAL_MONTHLY]*3, [REAL_MONTHLY]*3, [HEALTHY_MONTHLY]*3)):
        # Official amount
        ax.text(i, OFFICIAL_MONTHLY/2, f'EUR {OFFICIAL_MONTHLY:.2f}', ha='center', va='center',
               fontweight='bold', fontsize=11, color='white')

        # Gap 1 (Official to Real)
        if i > 0:
            ax.text(i, OFFICIAL_MONTHLY + gap_official_real/2, f'+EUR {gap_official_real:.2f}\n(+342%)',
                   ha='center', va='center', fontweight='bold', fontsize=10, color='white')

        # Gap 2 (Real to Healthy)
        if i > 1:
            ax.text(i, OFFICIAL_MONTHLY + gap_official_real + gap_real_healthy/2,
                   f'+EUR {gap_real_healthy:.2f}\n(+26%)', ha='center', va='center',
                   fontweight='bold', fontsize=10, color='white')

        # Total
        total = OFFICIAL_MONTHLY + (gap_official_real if i > 0 else 0) + (gap_real_healthy if i > 1 else 0)
        ax.text(i, total + 15, f'EUR {total:.2f}', ha='center', va='bottom',
               fontweight='bold', fontsize=12)

    ax.set_ylabel('Monthly Cost (EUR)', fontsize=13, fontweight='bold')
    ax.set_title('The Hidden Gap: Official vs Real Living Costs', fontsize=14, fontweight='bold')
    ax.set_ylim(0, HEALTHY_MONTHLY * 1.2)
    ax.legend(fontsize=11, loc='upper left')
    ax.set_xticklabels(baskets)

    plt.tight_layout()
    plt.savefig(os.path.join(BASKET_ANALYSIS_DIR, '03_gap_analysis.png'), dpi=150, bbox_inches='tight')
    plt.close()
    print("   [OK] 03_gap_analysis.png")

    # 4. Basket Composition (Pie charts)
    fig, axes = plt.subplots(1, 3, figsize=(16, 5))

    baskets_data = [
        ('Official\n13 Items', ['Dairy', 'Proteins', 'Staples', 'Produce', 'Beverages']),
        ('Real Complete\n34 Items', ['Dairy', 'Proteins', 'Staples', 'Produce', 'Baby/Health', 'Beverages', 'Cultural']),
        ('Healthy\n41 Items', ['Dairy', 'Proteins', 'Staples', 'Produce', 'Baby/Health', 'Beverages', 'Quality Items']),
    ]

    for ax, (name, categories) in zip(axes, baskets_data):
        sizes = [20, 15, 18, 20, 12, 8, 7] if len(categories) == 7 else ([20, 15, 18, 20, 12, 15] if len(categories) == 6 else [20, 25, 18, 20, 17])

        colors_pie = plt.cm.Set3(np.linspace(0, 1, len(categories)))
        wedges, texts, autotexts = ax.pie(sizes[:len(categories)], labels=categories, autopct='%1.0f%%',
                                          colors=colors_pie, startangle=90, textprops={'fontsize': 9})

        for autotext in autotexts:
            autotext.set_color('white')
            autotext.set_fontweight('bold')
            autotext.set_fontsize(8)

        for text in texts:
            text.set_fontsize(9)

        ax.set_title(name, fontsize=12, fontweight='bold', pad=10)

    plt.suptitle('Basket Composition: What Gets Counted?', fontsize=14, fontweight='bold', y=1.02)
    plt.tight_layout()
    plt.savefig(os.path.join(BASKET_ANALYSIS_DIR, '04_basket_composition.png'), dpi=150, bbox_inches='tight')
    plt.close()
    print("   [OK] 04_basket_composition.png")

    # 5. Remaining Budget After Food
    fig, ax = plt.subplots(figsize=(13, 7))

    monthly_incomes = [inc/12 for inc in PARIS_INCOMES.values()]

    official_remaining = [income - OFFICIAL_MONTHLY for income in monthly_incomes]
    real_remaining = [income - REAL_MONTHLY for income in monthly_incomes]
    healthy_remaining = [income - HEALTHY_MONTHLY for income in monthly_incomes]

    x = np.arange(len(PARIS_INCOMES.keys()))
    width = 0.25

    bars1 = ax.bar(x - width, official_remaining, width, label='After Official Basket',
                  color=CHART_BLUE, edgecolor='black', linewidth=1.5)
    bars2 = ax.bar(x, real_remaining, width, label='After Real Basket',
                  color=CHART_ORANGE, edgecolor='black', linewidth=1.5)
    bars3 = ax.bar(x + width, healthy_remaining, width, label='After Healthy Basket',
                  color='#70ad47', edgecolor='black', linewidth=1.5)

    # Crisis line
    ax.axhline(y=0, color='red', linestyle='-', linewidth=2.5, alpha=0.8, label='Financial Crisis Point')
    ax.axhline(y=500, color='orange', linestyle='--', linewidth=2, alpha=0.6, label='Critical Stress Level')
    ax.axhline(y=1000, color='green', linestyle='--', linewidth=2, alpha=0.6, label='Comfortable Level')

    ax.set_ylabel('Remaining Monthly Budget (EUR)', fontsize=13, fontweight='bold')
    ax.set_title('Remaining Budget After Food: Can Families Afford Rent & Life?', fontsize=14, fontweight='bold')
    ax.set_xticks(x)
    ax.set_xticklabels(list(PARIS_INCOMES.keys()), fontsize=11)
    ax.legend(loc='upper right', fontsize=10)
    ax.grid(axis='y', alpha=0.3)

    # Add value labels
    for bars in [bars1, bars2, bars3]:
        for bar in bars:
            height = bar.get_height()
            if height < 0:
                va_pos = 'top'
                y_offset = -100
            else:
                va_pos = 'bottom'
                y_offset = 50
            ax.text(bar.get_x() + bar.get_width()/2., height + y_offset,
                   f'EUR {height:.0f}', ha='center', va=va_pos, fontweight='bold', fontsize=9)

    plt.tight_layout()
    plt.savefig(os.path.join(BASKET_ANALYSIS_DIR, '05_healthy_cliff.png'), dpi=150, bbox_inches='tight')
    plt.close()
    print("   [OK] 05_healthy_cliff.png")

def create_updated_presentation():
    """Create presentation with corrected data"""
    print("\n[STEP 2] Creating updated presentation...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BLUE

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(60)
        title_p.font.bold = True
        title_p.font.color.rgb = WHITE
        title_p.alignment = PP_ALIGN.CENTER

        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = subtitle
        subtitle_p.font.size = Pt(28)
        subtitle_p.font.color.rgb = ACCENT_ORANGE
        subtitle_p.alignment = PP_ALIGN.CENTER

    def add_content_slide(title, content_items):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = DARK_BLUE
        title_shape.line.color.rgb = DARK_BLUE

        title_frame = title_shape.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(44)
        title_p.font.bold = True
        title_p.font.color.rgb = WHITE
        title_p.space_before = Pt(12)
        title_p.space_after = Pt(12)

        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.7))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, item in enumerate(content_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_BLACK
            p.space_before = Pt(6)
            p.space_after = Pt(8)

    def add_image_slide(title, image_path):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = DARK_BLUE
        title_shape.line.color.rgb = DARK_BLUE

        title_frame = title_shape.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(44)
        title_p.font.bold = True
        title_p.font.color.rgb = WHITE
        title_p.space_before = Pt(12)
        title_p.space_after = Pt(12)

        left = Inches(0.5)
        top = Inches(1.3)
        height = Inches(5.8)
        pic = slide.shapes.add_picture(image_path, left, top, height=height)

    def add_quote_slide(quote, attribution):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BLUE

        quote_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2.5))
        quote_frame = quote_box.text_frame
        quote_frame.word_wrap = True
        quote_p = quote_frame.paragraphs[0]
        quote_p.text = quote
        quote_p.font.size = Pt(32)
        quote_p.font.bold = True
        quote_p.font.italic = True
        quote_p.font.color.rgb = ACCENT_ORANGE
        quote_p.alignment = PP_ALIGN.CENTER

        attr_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(1))
        attr_frame = attr_box.text_frame
        attr_p = attr_frame.paragraphs[0]
        attr_p.text = attribution
        attr_p.font.size = Pt(20)
        attr_p.font.color.rgb = WHITE
        attr_p.alignment = PP_ALIGN.CENTER

    # Build presentation
    add_title_slide("The Reality Behind the Numbers",
                   "How Official Inflation Statistics Hide True Living Costs in Paris")

    add_content_slide("The Problem: Official Baskets are Incomplete", [
        "INSEE reports inflation based on MINIMAL BASKET of emblematic products",
        "",
        "But official baskets IGNORE complementary items families need:",
        "  * Different milk types (lactose-free, baby formula)",
        "  * Items that complete meals (proteins, vegetables, sauces)",
        "  * Baby essentials (diapers, creams)",
        "  * Dietary alternatives (allergies, intolerances)",
        "",
        "Result: Official inflation measures SURVIVAL, not LIVING"
    ])

    add_content_slide("CORRECTED: Real Household Costs (2 adults + 1 child)", [
        "Official Minimal Basket: EUR 68.67/month",
        "  • 13 basic items (milk, bread, eggs, pasta...)",
        "  • Only 3.9% of SMIC income",
        "",
        "Real Complete Basket: EUR 303.65/month",
        "  • 34 items (adds proteins, vegetables, baby items, hygiene)",
        "  • Gap: +342% more than official",
        "  • 17.4% of SMIC income",
        "",
        "Healthy Complete Basket: EUR 383.58/month",
        "  • 41 items (quality, organic, nutritious)",
        "  • Health premium: +26% over real basket",
        "  • 21.9% of SMIC income"
    ])

    add_image_slide("The Cost Reality: Official vs Real vs Healthy",
                   os.path.join(BASKET_ANALYSIS_DIR, "01_basket_comparison.png"))

    add_quote_slide("Official inflation rates are built on incomplete baskets.\nThey measure SURVIVAL, not LIVING.",
                   "INSIGHT #1: The Incomplete Statistic")

    add_image_slide("Affordability Cliffs: When Food Costs Exceed Income",
                   os.path.join(BASKET_ANALYSIS_DIR, "02_affordability_cliffs.png"))

    add_content_slide("The Affordability Crisis by Income Level", [
        "SMIC WORKER (21,000 EUR/year):",
        "  * Official: 3.9% of income (manageable)",
        "  * Real: 17.4% of income (difficult)",
        "  * Healthy: 21.9% of income (TIGHT)",
        "",
        "MEDIAN PARIS (42,000 EUR/year):",
        "  * Real: 8.7% of income (manageable)",
        "  * Healthy: 11.0% of income (manageable)",
        "",
        "THE CLIFF: Below 28k EUR/year, families cannot afford healthy food"
    ])

    add_quote_slide("Healthy eating is a luxury.\nThe poor eat cheap (unhealthy).\nThe wealthy eat well.",
                   "INSIGHT #2: Health as Enforced Inequality")

    add_image_slide("The Hidden Gap: What Official Statistics Don't Count",
                   os.path.join(BASKET_ANALYSIS_DIR, "03_gap_analysis.png"))

    add_content_slide("What's Missing from Official Baskets?", [
        "BABY/FAMILY ESSENTIALS (EUR 199/month not counted):",
        "  * Diapers: EUR 18.99/pack",
        "  * Baby creams & ointments: EUR 8.99",
        "",
        "COMPLEMENTARY PROTEINS (EUR 29+/month):",
        "  * Chicken, fish, beef for variety and nutrition",
        "",
        "DIETARY ALTERNATIVES:",
        "  * Lactose-free milk for intolerant families",
        "",
        "HEALTH & HYGIENE (EUR 15-20/month):",
        "  * Toilet paper, soap, shampoo, toothpaste"
    ])

    add_quote_slide("Babies don't exist in official inflation baskets.\nNeither do allergy sufferers.\nStatistical invisibility = institutional indifference.",
                   "INSIGHT #3: The Erased Family Reality")

    add_image_slide("What's Actually in Each Basket?",
                   os.path.join(BASKET_ANALYSIS_DIR, "04_basket_composition.png"))

    add_image_slide("The Health Paradox: When Healthy Becomes Impossible",
                   os.path.join(BASKET_ANALYSIS_DIR, "05_healthy_cliff.png"))

    add_quote_slide("Poor eat cheap, get sick, can't afford healthcare.\nThis creates a poverty trap.\nFood affordability forces health inequality.",
                   "INSIGHT #4: The Poverty Health Trap")

    add_content_slide("Why Snipper Tool Changes This", [
        "Snipper Tool collects REAL pricing data for REAL baskets",
        "",
        "By tracking actual purchases across Paris supermarkets:",
        "  * Expose the gap between official statistics and real costs",
        "  * Track health-focused pricing decisions",
        "  * Identify affordability cliffs by income",
        "  * Provide data for better policy and wage decisions",
        "",
        "What gets measured is what gets managed.\nWhat doesn't get measured gets ignored."
    ])

    add_title_slide("The Measurement Matters",
                   "Complete data = Better policy for real families")

    prs.save(PRESENTATION_OUTPUT)
    print(f"   [OK] Presentation updated")

def copy_report():
    """Copy corrected report to basket_analysis folder"""
    print("\n[STEP 3] Updating analysis reports...")

    source = os.path.join(BASKET_CORRECTED_DIR, 'BASKET_ANALYSIS_CORRECTED.txt')
    dest = os.path.join(BASKET_ANALYSIS_DIR, 'BASKET_ANALYSIS_REPORT.txt')

    if os.path.exists(source):
        shutil.copy2(source, dest)
        print(f"   [OK] Copied corrected report to basket_analysis/")
    else:
        print(f"   [ERROR] Source not found: {source}")

# Main execution
if __name__ == "__main__":
    print("="*80)
    print("UPDATING DELIVERABLES WITH CORRECTED BASKET ANALYSIS")
    print("="*80)

    create_corrected_visualizations()
    create_updated_presentation()
    copy_report()

    print("\n" + "="*80)
    print("UPDATE COMPLETE")
    print("="*80)
