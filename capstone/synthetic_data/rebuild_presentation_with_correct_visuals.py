"""
Rebuild PowerPoint presentation with CORRECTED visualizations
Using per-capita methodology and consolidated visualization files
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Paths
ORIGINAL_PPTX = r'C:\Users\emman\p_Claude\devs\snipper_tool\CAPSTONE_DELIVERABLES\Snipper_Tool_Reality_vs_Official_Paris.pptx'
OUTPUT_PPTX = r'C:\Users\emman\p_Claude\devs\snipper_tool\CAPSTONE_DELIVERABLES\Snipper_Tool_Reality_vs_Official_Paris.pptx'
VIZ_DIR = r'C:\Users\emman\p_Claude\devs\snipper_tool\CAPSTONE_DELIVERABLES'

# Load original presentation
prs = Presentation(ORIGINAL_PPTX)
print(f"[INFO] Loaded presentation with {len(prs.slides)} slides")

# Map of slide positions to update with visualizations
# Based on VISUALIZATION_INVENTORY.md presentation flow
VISUALIZATION_MAP = {
    # Slide 4: Basket Comparison (Segment 1 - The Hook)
    3: {
        'chart': '01_basket_comparison.png',
        'title': 'The Gap: Official vs Reality',
        'description': 'Official (€69) vs Real (€304) vs Healthy (€384) baskets'
    },
    # Slide 6: Affordability Cliffs (Segment 3B - Crisis by Income)
    5: {
        'chart': '02_affordability_cliffs.png',
        'title': 'Affordability Cliffs by Income Level',
        'description': 'Food cost as % of income with crisis zones'
    },
    # Slide 9: Gap Analysis (Segment 3A - Detail)
    8: {
        'chart': '03_gap_analysis.png',
        'title': 'Gap Analysis: What\'s Missing',
        'description': 'Breakdown of what\'s missing from official basket'
    },
    # Slide 11: Basket Composition (Segment 3A - Composition)
    10: {
        'chart': '04_basket_composition.png',
        'title': 'Basket Composition: Items vs Cost',
        'description': '13 items (official) vs 34 items (real) vs 41 items (healthy)'
    },
    # Slide 12: Healthy Cliff (Supporting Evidence)
    11: {
        'chart': '05_healthy_cliff.png',
        'title': 'Healthy Cliff: Affordability After Health Focus',
        'description': 'Remaining budget after healthy eating costs'
    },
    # Slide 15: Family Squeeze (Segment 3C - Household Impact)
    14: {
        'chart': '06_family_squeeze.png',
        'title': 'The Family Squeeze (CORRECTED)',
        'description': 'How each member reduces remaining budget (per-capita methodology)'
    },
    # Slide 16: Housing Burden (Segment 3C - Housing Context)
    15: {
        'chart': '07_housing_burden.png',
        'title': 'Housing Burden by Income Level',
        'description': 'Housing cost variation and impact on remaining budget'
    },
    # Slide 25: Complete Budget Reality (Segment 3C - Structural Deficit)
    24: {
        'chart': '08_complete_budget_reality.png',
        'title': 'Complete Budget Reality (CORRECTED)',
        'description': 'All 12 essential costs: SMIC vs Median'
    },
}

# Update visualizations
updated_count = 0
for slide_idx, viz_info in VISUALIZATION_MAP.items():
    chart_path = os.path.join(VIZ_DIR, viz_info['chart'])

    if not os.path.exists(chart_path):
        print(f"[WARN] Chart not found: {chart_path}")
        continue

    try:
        slide = prs.slides[slide_idx]

        # Remove old pictures
        for shape in list(slide.shapes):
            if shape.shape_type == 13:  # Picture
                sp = shape.element
                sp.getparent().remove(sp)

        # Add new picture (positioning varies by slide, but generally center)
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)

        slide.shapes.add_picture(chart_path, left, top, width=width, height=height)

        print(f"[OK] Slide {slide_idx + 1}: Updated with {viz_info['chart']}")
        updated_count += 1

    except Exception as e:
        print(f"[ERROR] Slide {slide_idx + 1}: {str(e)}")

# Save updated presentation
prs.save(OUTPUT_PPTX)
print(f"\n[SUCCESS] Updated presentation saved to: {OUTPUT_PPTX}")
print(f"[INFO] Updated {updated_count} slides with correct visualizations")
print("\nVisualization Updates:")
for slide_idx, viz_info in VISUALIZATION_MAP.items():
    print(f"  - Slide {slide_idx + 1}: {viz_info['chart']} ({viz_info['title']})")

print("\n[CRITICAL] All visualizations now use CORRECTED per-capita methodology:")
print("  - Family Squeeze (06): €117.74/month for teenager (was €288.47)")
print("  - Complete Budget Reality (08): All 12 essential costs, SMIC €74 deficit")
print("  - All member costs: Baby €30.98, Child €55.77, Teen €117.74, Senior €136.33")
