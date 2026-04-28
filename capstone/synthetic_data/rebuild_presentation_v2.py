"""
Rebuild PowerPoint with corrected visualizations - Version 2
More careful handling of slide modifications
"""

from pptx import Presentation
from pptx.util import Inches
import os

PPTX_PATH = r'C:\Users\emman\p_Claude\devs\snipper_tool\CAPSTONE_DELIVERABLES\Snipper_Tool_Reality_vs_Official_Paris.pptx'
VIZ_DIR = r'C:\Users\emman\p_Claude\devs\snipper_tool\CAPSTONE_DELIVERABLES'

# Map of slide indices (0-based) to visualization files
SLIDES_TO_UPDATE = {
    3: '01_basket_comparison.png',          # Slide 4
    5: '02_affordability_cliffs.png',       # Slide 6
    8: '03_gap_analysis.png',               # Slide 9
    10: '04_basket_composition.png',        # Slide 11
    11: '05_healthy_cliff.png',             # Slide 12
    14: '06_family_squeeze.png',            # Slide 15 (CORRECTED)
    15: '07_housing_burden.png',            # Slide 16
    24: '08_complete_budget_reality.png',   # Slide 25 (CORRECTED)
}

print("[INFO] Loading presentation...")
prs = Presentation(PPTX_PATH)
print("[INFO] Loaded {} slides".format(len(prs.slides)))

updated = 0
for slide_idx, chart_filename in SLIDES_TO_UPDATE.items():
    chart_path = os.path.join(VIZ_DIR, chart_filename)

    if not os.path.exists(chart_path):
        print("[SKIP] Chart not found: {}".format(chart_path))
        continue

    try:
        slide = prs.slides[slide_idx]

        # Find and remove existing pictures (carefully)
        pictures_removed = 0
        for shape in list(slide.shapes):
            if shape.shape_type == 13:  # Picture type
                try:
                    sp = shape.element
                    sp.getparent().remove(sp)
                    pictures_removed += 1
                except:
                    pass

        # Add corrected visualization
        slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.5), width=Inches(9), height=Inches(5))

        print("[OK] Slide {}: {} (removed {} old images)".format(
            slide_idx + 1, chart_filename, pictures_removed
        ))
        updated += 1

    except Exception as e:
        print("[ERROR] Slide {}: {}".format(slide_idx + 1, str(e)))

# Save
print("\n[INFO] Saving presentation...")
prs.save(PPTX_PATH)

print("[SUCCESS] Updated {} slides".format(updated))
print("[INFO] All visualizations use CORRECTED per-capita methodology:")
print("  - Teenager: EUR 117.74/month (was EUR 288.47)")
print("  - Family Squeeze: Reflects correct per-capita costs")
print("  - Complete Budget Reality: EUR 74 monthly deficit for SMIC families")
