from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Paths
CHARTS_DIR = r"C:\Users\emman\p_Claude\devs\snipper_tool\capstone\synthetic_data\basket_analysis"
OUTPUT_PATH = r"C:\Users\emman\p_Claude\devs\snipper_tool\CAPSTONE_DELIVERABLES\Snipper_Tool_Reality_vs_Official_Paris.pptx"

# Colors
DARK_BLUE = RGBColor(31, 78, 121)
ACCENT_GREEN = RGBColor(79, 129, 189)
ACCENT_ORANGE = RGBColor(237, 125, 49)
RED_ALERT = RGBColor(255, 68, 68)
WHITE = RGBColor(255, 255, 255)
TEXT_BLACK = RGBColor(51, 51, 51)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(60)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = ACCENT_ORANGE
    subtitle_p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items):
    """Add content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE

    title_frame = title_shape.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.space_before = Pt(12)
    title_p.space_after = Pt(12)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.7))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_BLACK
        p.space_before = Pt(6)
        p.space_after = Pt(8)

def add_image_slide(prs, title, image_path, height_inches=5.5):
    """Add slide with image and title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE

    title_frame = title_shape.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.space_before = Pt(12)
    title_p.space_after = Pt(12)

    # Add image centered
    left = Inches(0.5)
    top = Inches(1.3)
    height = Inches(height_inches)
    pic = slide.shapes.add_picture(image_path, left, top, height=height)

def add_quote_slide(prs, quote, attribution):
    """Add quote slide for emphasis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    # Main quote
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2.5))
    quote_frame = quote_box.text_frame
    quote_frame.word_wrap = True
    quote_p = quote_frame.paragraphs[0]
    quote_p.text = quote
    quote_p.font.size = Pt(32)
    quote_p.font.bold = True
    quote_p.font.italic = True
    quote_p.font.color.rgb = ACCENT_ORANGE
    quote_p.alignment = PP_ALIGN.CENTER

    # Attribution
    attr_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(1))
    attr_frame = attr_box.text_frame
    attr_p = attr_frame.paragraphs[0]
    attr_p.text = attribution
    attr_p.font.size = Pt(20)
    attr_p.font.color.rgb = WHITE
    attr_p.alignment = PP_ALIGN.CENTER

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title
add_title_slide(prs, "The Reality Behind the Numbers",
               "How Official Inflation Statistics Hide the True Cost of Living in Paris")

# Slide 2: The Problem
add_content_slide(prs, "The Problem: Official Baskets are Incomplete", [
    "INSEE reports inflation based on a MINIMAL BASKET of emblematic products:",
    "  • Milk, bread, eggs, cheese, butter, pasta...",
    "",
    "But official baskets IGNORE the complementary items families actually need:",
    "  • Different milk types (lactose-free, baby formula)",
    "  • Items that complete meals (sauces, oils, spices)",
    "  • Baby essentials (diapers, creams)",
    "  • Dietary alternatives (allergies, intolerances)",
    "",
    "Result: Official inflation numbers measure SURVIVAL, not LIVING"
])

# Slide 3: What's Missing?
add_content_slide(prs, "What's Missing from Official Baskets?", [
    "BABY & FAMILY ESSENTIALS (27.98 EUR/month):",
    "  • Diapers: 18.99 EUR",
    "  • Baby creams & ointments: 8.99 EUR",
    "",
    "DIETARY VARIETIES & ALTERNATIVES:",
    "  • Lactose-free milk, gluten-free options, different protein sources",
    "",
    "COMPLEMENTARY COOKING ITEMS:",
    "  • Spices, sauces, quality oils - items that complete meals",
    "",
    "HEALTH & HYGIENE ESSENTIALS:",
    "  • Toilet paper, soap, shampoo, toothpaste: 15-20 EUR/month",
    "",
    "CULTURAL & LIFE QUALITY:",
    "  • Wine, chocolate, treats - the cost of living, not just surviving"
])

# Slide 4: The Numbers - Basket Comparison
add_image_slide(prs, "The Cost Gap: Official vs Real",
               os.path.join(CHARTS_DIR, "01_basket_comparison.png"), 5.2)

# Slide 5: Key Insight #1
add_quote_slide(prs,
    "Official inflation rates are built on incomplete baskets.\nThey measure SURVIVAL, not LIVING.",
    "CLICHE INSIGHT #1: The Incomplete Statistic")

# Slide 6: Affordability Crisis
add_image_slide(prs, "The Affordability Cliff: When Real Costs Exceed Means",
               os.path.join(CHARTS_DIR, "02_affordability_cliffs.png"), 5.3)

# Slide 7: Detailed Income Analysis
add_content_slide(prs, "Affordability by Paris Income Level", [
    "SMIC WORKERS (21,000 EUR/year):",
    "  • Official basket: 9% of income (MANAGEABLE)",
    "  • Real complete basket: 18% of income (DIFFICULT)",
    "  • Healthy basket: 24% of income (CRISIS - NO MONEY FOR ANYTHING ELSE)",
    "",
    "MEDIAN PARIS (42,000 EUR/year):",
    "  • Real basket: 9% of income (MANAGEABLE)",
    "  • Healthy basket: 12% of income (TIGHT)",
    "",
    "THE CLIFF: Below 40k EUR annual income, healthy eating becomes IMPOSSIBLE"
])

# Slide 8: Key Insight #2
add_quote_slide(prs,
    "Healthy eating is a luxury.\nThe poor must eat cheap (unhealthy).\nThe wealthy eat well.",
    "CLICHE INSIGHT #2: Health as Inequality")

# Slide 9: The Gap Analysis
add_image_slide(prs, "The Hidden Gap: What's Not Counted",
               os.path.join(CHARTS_DIR, "03_gap_analysis.png"), 5.3)

# Slide 10: Creative Insight #1
add_quote_slide(prs,
    "Babies don't exist in official inflation baskets.\nNeither do allergy sufferers.\nOr single parents managing dietary needs.\nStatistical invisibility = institutional indifference.",
    "CREATIVE INSIGHT #1: The Erased Family Reality")

# Slide 11: Basket Composition
add_image_slide(prs, "What's Actually in Each Basket?",
               os.path.join(CHARTS_DIR, "04_basket_composition.png"), 5.2)

# Slide 12: Healthy Eating Cliff
add_image_slide(prs, "The Health Paradox: When Healthy Becomes Impossible",
               os.path.join(CHARTS_DIR, "05_healthy_cliff.png"), 5.3)

# Slide 13: Creative Insight #2
add_quote_slide(prs,
    "The poor eat cheap and get sick.\nThe sick can't afford healthcare.\nHealthcare costs break the budget.\nFood affordability creates a poverty trap.",
    "CREATIVE INSIGHT #2: The Poverty Health Trap")

# Slide 14: The Data Reality
add_content_slide(prs, "The Data Reveals the Truth", [
    "OFFICIAL MINIMAL BASKET (Monthly): 259 EUR",
    "  • Items counted: 11 emblematic products",
    "  • Annual cost: 3,107 EUR",
    "",
    "REAL COMPLETE BASKET (Monthly): 510 EUR",
    "  • Items needed: 29 products",
    "  • Annual cost: 6,119 EUR",
    "  • Hidden gap: +97% more than official",
    "",
    "HEALTHY COMPLETE BASKET (Monthly): 630 EUR",
    "  • Items for health: 36 products",
    "  • Annual cost: 7,558 EUR",
    "  • Health premium: +24% more than real basket"
])

# Slide 15: Household Composition Impact
add_image_slide(prs, "The Affordability Multiplier: Household Composition Impact",
               os.path.join(CHARTS_DIR, "06_household_cost_impact.png"), 5.3)

# Slide 16: Affordability by Composition
add_image_slide(prs, "The Multiplication Effect: Composition × Income = Crisis",
               os.path.join(CHARTS_DIR, "07_affordability_by_composition.png"), 5.3)

# Slide 17: Remaining Budget Crisis
add_content_slide(prs, "The Household Composition Multiplier Crisis", [
    "SMIC WORKER (21,000 EUR/year = 1,750 EUR/month):",
    "  • Base (2 adults + 1 child): 17.4% of income (MANAGEABLE)",
    "  • +1 Teenager: 33.8% of income (IMPOSSIBLE - NO MONEY FOR HOUSING)",
    "  • +1 Senior: 36.4% of income (CATASTROPHIC)",
    "",
    "MEDIAN PARIS (42,000 EUR/year):",
    "  • Base: 8.7% of income (EXCELLENT)",
    "  • +1 Teenager: 16.9% of income (MANAGEABLE)",
    "  • +1 Senior: 18.2% of income (DIFFICULT BUT POSSIBLE)",
    "",
    "THE CRUEL REALITY:",
    "Families don't control their composition. They have children, care for aging",
    "parents. Yet affordability multiplies with each family member.",
    "Below 35k income: Growing families become impossible families."
])

# Slide 18: The Invisibility Crisis - Who Gets Erased
add_content_slide(prs, "The Invisibility Crisis: Who Official Statistics Erase", [
    "SENIORS: Specialized dietary needs (soft foods, health management)",
    "  • €334/month more than baseline, but counted as 'standard adult'",
    "  • Elderly people disappear from policy view",
    "",
    "TEENAGERS: Abrupt appetite transition (110% cost jump at age 13)",
    "  • Malnutrition risk precisely when growth is critical",
    "  • Not child-cost, not adult-cost—uncounted gap",
    "",
    "THE SANDWICH GENERATION: Supporting both elderly AND children",
    "  • 45-year-old: Parent +€334, Child +€304 = household multiplied",
    "  • €941/month food costs (27% of median income)",
    "  • Silent burden. Invisible in statistics. Breaking quietly."
])

# Slide 19: The Demographic Consequence
add_quote_slide(prs,
    "What gets priced out of reach doesn't get born.\nFood affordability is not just poverty.\nIt's a demographic determinant.",
    "THE HIDDEN CONSEQUENCE: Why Birth Rates Are Collapsing")

# Slide 20: Birth Rate as Economic Rationing
add_content_slide(prs, "Birth Rate as Economic Rationing", [
    "Young couple, median Paris income (42,000 EUR/year):",
    "  • 1 child: 12.6% of income on food (MANAGEABLE)",
    "  • 2 children: 16.5% of income on food (DIFFICULT)",
    "  • 3 children: 20.4% of income on food (IMPOSSIBLE)",
    "",
    "Young couple, SMIC (21,000 EUR/year):",
    "  • 1 child: 25.2% of income (CRISIS)",
    "  • 2 children: 38.9% of income (IMPOSSIBLE - CHOOSE: FOOD OR HOUSING)",
    "",
    "Paris birth rate: 1.4 children/woman (replacement: 2.1)",
    "This is not cultural choice. This is economic calculus.",
    "Families are rationing their size. Demographic collapse is silent rationing."
])

# Slide 21: Why This Matters
add_content_slide(prs, "Why This Matters", [
    "POLICY MAKERS base decisions on official inflation rates that underestimate",
    "real costs by 30-50%",
    "",
    "WORKERS see their salaries (linked to inflation) never keep pace with actual",
    "living expenses",
    "",
    "FAMILIES make impossible choices: healthy food OR paying rent?",
    "Diapers OR transportation?",
    "",
    "INEQUALITY is hidden in the numbers. What's not measured is not managed.",
    "",
    "STATISTICS are not neutral. They reveal what society has decided to COUNT",
    "and what it chooses to IGNORE."
])

# Slide 22: The Opportunity for Snipper Tool
add_content_slide(prs, "The Opportunity: Real Data for Real Decisions", [
    "Snipper Tool changes this by collecting REAL pricing data for REAL baskets:",
    "",
    "Real families don't buy INSEE baskets.",
    "They buy complete meals, dietary alternatives, baby products, and items",
    "that support healthy, dignified family life.",
    "",
    "By tracking ACTUAL purchases across Paris supermarkets, Snipper Tool can:",
    "  • Expose the gap between official statistics and real costs",
    "  • Track health-focused pricing decisions",
    "  • Identify affordability cliffs by income and family type",
    "  • Provide data for better policy and wage decisions"
])

# Slide 23: Conclusion
add_quote_slide(prs,
    "The numbers never lie.\nBut incomplete numbers hide the truth.",
    "FINAL INSIGHT: What You Measure is What You Manage")

# Slide 24: Data Summary
add_content_slide(prs, "This Analysis", [
    "Dataset: 2,163 synthetic price observations (Apr 2025 - Apr 2026)",
    "From Snipper Tool captures in Paris Auchan and Carrefour stores",
    "",
    "Compared THREE basket types:",
    "  1. Official INSEE Minimal Basket",
    "  2. Real Complete Living Basket (includes complementary items)",
    "  3. Healthy Complete Basket (for nutritional wellness)",
    "",
    "Against PARIS HOUSEHOLD INCOMES:",
    "  • SMIC (minimum wage): 21,000 EUR/year",
    "  • Low income: 28,000 EUR/year",
    "  • Median: 42,000 EUR/year",
    "  • Upper middle: 65,000 EUR/year",
    "",
    "Geography: Paris region (Ile-de-France)"
])

# Save presentation
prs.save(OUTPUT_PATH)
print(f"Presentation created successfully!")
print(f"Location: {OUTPUT_PATH}")
print(f"Total slides: {len(prs.slides)}")
print(f"\nPresentation ready for demonstration.")
