"""
Rebuild PowerPoint with corrected visualizations - Version 3
Uses backup and very careful image handling
"""

from pptx import Presentation
from pptx.util import Inches
import os
import shutil
import tempfile

PPTX_PATH = r'C:\Users\emman\p_Claude\devs\snipper_tool\CAPSTONE_DELIVERABLES\Snipper_Tool_Reality_vs_Official_Paris.pptx'
BACKUP_PATH = PPTX_PATH + '.backup'
VIZ_DIR = r'C:\Users\emman\p_Claude\devs\snipper_tool\CAPSTONE_DELIVERABLES'

# Map of slide indices (0-based) to visualization files
SLIDES_TO_UPDATE = {
    3: '01_basket_comparison.png',          # Slide 4
    5: '02_affordability_cliffs.png',       # Slide 6
    8: '03_gap_analysis.png',               # Slide 9
    10: '04_basket_composition.png',        # Slide 11
    11: '05_healthy_cliff.png',             # Slide 12
    14: '06_family_squeeze.png',            # Slide 15 (CORRECTED)
    15: '07_housing_burden.png',            # Slide 16
    24: '08_complete_budget_reality.png',   # Slide 25 (CORRECTED)
}

print("[INFO] Creating backup...")
if not os.path.exists(BACKUP_PATH):
    shutil.copy2(PPTX_PATH, BACKUP_PATH)
    print("[OK] Backup created: {}".format(BACKUP_PATH))

print("[INFO] Loading presentation...")
prs = Presentation(PPTX_PATH)
print("[INFO] Loaded {} slides".format(len(prs.slides)))

updated = 0
failed = 0

for slide_idx, chart_filename in SLIDES_TO_UPDATE.items():
    chart_path = os.path.join(VIZ_DIR, chart_filename)

    if not os.path.exists(chart_path):
        print("[SKIP] Chart not found: {}".format(chart_path))
        continue

    try:
        slide = prs.slides[slide_idx]

        # Strategy: Add new image first, then remove old ones
        # This is safer than removing first
        pic = slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.5), width=Inches(9), height=Inches(5))

        # Now remove old pictures (newer python-pptx handles this better after adding)
        old_count = 0
        for shape in list(slide.shapes):
            if shape.shape_type == 13 and shape != pic:  # Picture type, but not the one we just added
                try:
                    sp = shape.element
                    sp.getparent().remove(sp)
                    old_count += 1
                except:
                    pass

        print("[OK] Slide {}: {} (replaced {} old)".format(
            slide_idx + 1, chart_filename, old_count
        ))
        updated += 1

    except Exception as e:
        print("[ERROR] Slide {}: {}".format(slide_idx + 1, str(e)))
        failed += 1

print("\n[INFO] Saving presentation...")
try:
    # Use temp file first to ensure atomic write
    temp_fd, temp_path = tempfile.mkstemp(suffix='.pptx')
    os.close(temp_fd)

    prs.save(temp_path)

    # Verify temp file
    test_prs = Presentation(temp_path)
    print("[OK] Temp file verified")

    # Move to final location
    shutil.move(temp_path, PPTX_PATH)
    print("[SUCCESS] Presentation saved")

except Exception as e:
    print("[ERROR] Save failed: {}".format(str(e)))
    print("[INFO] Restoring from backup...")
    if os.path.exists(BACKUP_PATH):
        shutil.copy2(BACKUP_PATH, PPTX_PATH)

print("\n[SUMMARY]")
print("  Updated: {} slides".format(updated))
print("  Failed: {} slides".format(failed))
print("  Backup: {}".format(BACKUP_PATH))
print("\n[INFO] All visualizations use CORRECTED per-capita methodology:")
print("  - Teenager: EUR 117.74/month (was EUR 288.47)")
print("  - Family Squeeze: Slide 15 - Corrected remaining budgets")
print("  - Complete Budget Reality: Slide 25 - EUR 74 deficit")
