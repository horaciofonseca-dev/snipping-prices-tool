from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BLUE = RGBColor(31, 78, 121)
ACCENT_GREEN = RGBColor(79, 129, 189)
LIGHT_GRAY = RGBColor(242, 242, 242)
WHITE = RGBColor(255, 255, 255)
TEXT_BLACK = RGBColor(51, 51, 51)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(66)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(32)
    subtitle_p.font.color.rgb = ACCENT_GREEN
    subtitle_p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE

    title_frame = title_shape.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.space_before = Pt(10)
    title_p.space_after = Pt(10)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = item
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_BLACK
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.level = 0

    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE

    title_frame = title_shape.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.space_before = Pt(10)
    title_p.space_after = Pt(10)

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4.3), Inches(6))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    # Left title
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(12)

    # Left items
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(19)
        p.font.color.rgb = TEXT_BLACK
        p.space_after = Pt(10)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4.3), Inches(6))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    # Right title
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(12)

    # Right items
    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(19)
        p.font.color.rgb = TEXT_BLACK
        p.space_after = Pt(10)

    return slide

# Slide 1: Title
add_title_slide(prs, "Snipper Tool", "Data-as-a-Service Platform\nMarket Intelligence & Pricing Analytics")

# Slide 2: The Problem
add_content_slide(prs, "The Market Problem", [
    "Market research is SLOW: Manual data collection takes 15-20 hours per dataset",
    "Web scraping is RISKY: Blocked by anti-bot measures, legal/GDPR concerns",
    "Market data is EXPENSIVE: Analyst fees + licensing costs €500-1,000+ per month",
    "Data is STALE: By the time insights arrive, markets have moved",
    "SOLUTION: Semi-automatic visual intelligence - fast, accurate, legal & scalable"
])

# Slide 3: How it Works
add_content_slide(prs, "How Snipper Tool Works: 3-Step Process", [
    "1. CAPTURE: User takes screenshot with smart overlay (Alt+C hotkey)",
    "2. DETECT: AI-powered OCR automatically extracts prices & product data",
    "3. ANALYZE: Data exported as structured JSON for immediate analysis",
    "",
    "TIME: 40 seconds per product vs 2 minutes manual (5x faster)",
    "ACCURACY: 100% validated data ready for analytics & reporting"
])

# Slide 4: Financial Case
add_content_slide(prs, "Financial Case: Supermarket Monitoring", [
    "Weekly Collection: 500 products across 5 stores (€146.74/week = €7,630/year)",
    "Manual Alternative: Would cost €45,500/year (€875/week)",
    "ANNUAL SAVINGS: €37,870 per customer (83% cost reduction)",
    "",
    "YOUR REVENUE: Sell for €200-400/month = €2,400-4,800/year per customer",
    "Gross margin: 68-84% (You get €1,600-3,000 of the €7,630 savings)"
])

# Slide 5: Revenue Model
add_content_slide(prs, "3-Stream Revenue Model", [
    "Stream 1: SaaS Subscriptions",
    "  € €50-200/month: Solo analyst (1 store, 100 products/week)",
    "  € €200-500/month: Enterprise tier (5+ stores, 1000+ products/week)",
    "",
    "Stream 2: Data Licensing - €500-5,000/month for aggregated datasets",
    "Stream 3: Consulting - €5,000-50,000 for custom deployments"
])

# Slide 6: Application 1 - Retail
add_two_column_slide(prs, "Application #1: Retail Price Monitoring (Current)",
    "Market Opportunity",
    [
        "Target: Supermarket chains, CPG brands, price platforms",
        "Market: €200M+ in retail analytics annually",
        "Uses: Competitive pricing, promotions, margins",
        "Geography: EU-wide (France, Spain, Germany, UK)"
    ],
    "Revenue Potential",
    [
        "50 customers @ €250/month = €150,000/year",
        "Data licensing: €2K-5K/month = €24K-60K/year",
        "Consulting: 10 x €10K = €100K/year",
        "TOTAL: €274K-310K Year 1"
    ])

# Slide 7: Application 2 - Real Estate
add_two_column_slide(prs, "Application #2: Real Estate Valuation",
    "What Changes",
    [
        "OCR extracts: Property prices, location, features, sqft",
        "Targets: Zillow, Rightmove, Immobilienscout24",
        "Data points: 20+ property attributes per capture",
        "Effort: 40% new code (same core platform)"
    ],
    "Revenue Potential",
    [
        "Market 3x larger than retail",
        "100 customers @ €400/month = €480K/year",
        "Valuation reports: €2K-10K per engagement",
        "TOTAL: €600K-800K Year 1"
    ])

# Slide 8: Application 3 - Automotive
add_two_column_slide(prs, "Application #3: Vehicle Pricing",
    "What Changes",
    [
        "OCR extracts: VIN, mileage, price, condition, specs",
        "Targets: Autotrader, Cars.com, Edmunds, dealers",
        "Data points: 30+ vehicle attributes per capture",
        "Effort: 35% new code (inventory tracking module)"
    ],
    "Revenue Potential",
    [
        "Used car market: €150B+ annually in EU",
        "120 customers @ €350/month = €504K/year",
        "Market reports: €5K-50K per deal",
        "TOTAL: €650K-1M Year 1"
    ])

# Slide 9: Application 4 - Hospitality
add_two_column_slide(prs, "Application #4: Hotel Rate Monitoring",
    "What Changes",
    [
        "OCR extracts: Room type, nightly rate, occupancy",
        "Targets: Booking.com, Hotels.com, Expedia",
        "Data points: 20+ hotel attributes per capture",
        "Effort: 40% new code (calendar tracking)"
    ],
    "Revenue Potential",
    [
        "Hospitality market highly price-sensitive",
        "80 customers @ €300/month = €288K/year",
        "Revenue optimization: €10K-100K per chain",
        "TOTAL: €450K-700K Year 1"
    ])

# Slide 10: Scalability
add_content_slide(prs, "The Scalability Advantage", [
    "REUSABLE CORE: Screenshot overlay + OCR + metadata works across ALL industries",
    "FASTER EXPANSION: Each new vertical takes 4-6 weeks (30-40% code changes)",
    "MARGINAL COSTS: Infrastructure scales logarithmically, not linearly",
    "",
    "Year 1 Revenue Path:",
    "  Q1-Q2: Retail (€150K) | Q3: Real Estate (€250K) | Q4: Auto (€400K) | Year 2: €1-2M"
])

# Slide 11: GTM Strategy
add_content_slide(prs, "Go-to-Market Strategy", [
    "Phase 1 (Months 1-3): Refine retail, land 10-15 customers",
    "  => Prove unit economics, build case studies",
    "",
    "Phase 2 (Months 4-6): Launch real estate, expand retail to 30+ customers",
    "  => Bundle offer: Retail + Real Estate (25% discount lock-in)",
    "",
    "Phase 3 (Months 7-12): Add automotive, target SMBs & mid-market",
    "  => Multi-industry customers get enterprise pricing (€5K+/month)",
    "",
    "Phase 4 (Year 2): Launch data licensing & analytics SaaS"
])

# Slide 12: Competitive Advantages
add_content_slide(prs, "Why Snipper Tool Wins", [
    "Speed: 5x faster than manual (5.07h vs 19h/week)",
    "Accuracy: 100% validated data (vs 85% manual, 10% scraping)",
    "Legal: Zero IP bans, no GDPR violations (vs scraping risks)",
    "Cost: €7,630/year operational (vs €45.5K manual or €235K scraping)",
    "Scalable: Add stores/products without adding people",
    "Repeatable: Works week after week with zero maintenance"
])

# Slide 13: Financial Projections
add_content_slide(prs, "5-Year Financial Projections", [
    "Year 1: €275K revenue | €175K profit (64% margin)",
    "Year 2: €850K revenue | €600K profit (71% margin)",
    "Year 3: €2.1M revenue | €1.65M profit (79% margin)",
    "Year 4: €4M revenue | €3.3M profit (83% margin)",
    "Year 5: €8M+ revenue | 85%+ margin (data licensing focus)"
])

# Slide 14: Implementation Roadmap
add_content_slide(prs, "6-Month Implementation Roadmap", [
    "Month 1: Polish retail MVP, prepare for beta (2-3 pilot customers)",
    "Month 2: Real estate module dev + 5-10 retail customers",
    "Month 3: Launch real estate beta + 20 retail customers",
    "Month 4: Automotive module dev + Enterprise tier (€5K+/month)",
    "Month 5: Automotive launch + Hospitality dev + multi-industry customers",
    "Month 6: 3 verticals live, 50+ customers, €300K+ MRR target"
])

# Slide 15: Closing
add_title_slide(prs, "Transform Data Collection into Revenue",
    "From €7,630 cost savings to\n€8M+ recurring revenue in 5 years")

# Save presentation
output_path = r"C:\Users\emman\p_Claude\devs\snipper_tool\Snipper_Tool_DaaS_Pitch.pptx"
prs.save(output_path)
print(f"PowerPoint created successfully!")
print(f"Location: {output_path}")
print(f"Total slides: {len(prs.slides)}")
